"""
============================================================================
CHURN CAUSALITY - EXECUTIVE PRESENTATION GENERATOR
============================================================================
Purpose: Create management-ready PowerPoint from causality analysis
Input: Charts and data from churn causality analysis
Output: Executive presentation with causal insights
============================================================================
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

print("=" * 80)
print("CHURN CAUSALITY - EXECUTIVE PRESENTATION BUILDER")
print("=" * 80)
print()

# ============================================================================
# CONFIGURATION
# ============================================================================

CHART_DIR = "/dbfs/FileStore/example_charts"
DATA_DIR = "/dbfs/FileStore/example_data"
OUTPUT_DIR = "/dbfs/FileStore/presentations"

os.makedirs(OUTPUT_DIR, exist_ok=True)

# Corporate colors
class Colors:
    PRIMARY = RGBColor(31, 71, 136)
    DANGER = RGBColor(231, 76, 60)
    WARNING = RGBColor(243, 156, 18)
    SUCCESS = RGBColor(39, 174, 96)
    INFO = RGBColor(52, 152, 219)
    SECONDARY = RGBColor(149, 165, 166)
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(44, 62, 80)
    LIGHT_BG = RGBColor(236, 240, 241)

colors = Colors()

# ============================================================================
# LOAD DATA
# ============================================================================

print("1. Loading causality analysis data...")

try:
    feature_importance = pd.read_csv(f"{DATA_DIR}/churn_feature_importance.csv")
    category_impact = pd.read_csv(f"{DATA_DIR}/churn_category_impact.csv")
    insights = pd.read_csv(f"{DATA_DIR}/churn_executive_insights.csv")
    
    print("   ✓ All data loaded successfully")
    print()
except Exception as e:
    print(f"   ✗ Error: {e}")
    print("   Run churn_causality_analysis.py first!")
    exit(1)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Create branded title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors.PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = colors.INFO
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide_with_chart(prs, title, chart_path, notes=None):
    """Add slide with chart and optional notes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors.PRIMARY
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    
    # Chart
    if os.path.exists(chart_path):
        if notes:
            # Chart on left, notes on right
            slide.shapes.add_picture(chart_path, Inches(0.3), Inches(1.1), width=Inches(6))
            
            # Notes box
            notes_box = slide.shapes.add_shape(
                1, Inches(6.5), Inches(1.3), Inches(3.2), Inches(3.8)
            )
            notes_box.fill.solid()
            notes_box.fill.fore_color.rgb = colors.LIGHT_BG
            notes_box.line.color.rgb = colors.PRIMARY
            notes_box.line.width = Pt(2)
            
            # Notes header
            notes_title = slide.shapes.add_textbox(Inches(6.7), Inches(1.5), Inches(2.8), Inches(0.4))
            notes_title.text = "💡 Key Insights"
            p = notes_title.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = colors.PRIMARY
            
            # Notes content
            notes_content = slide.shapes.add_textbox(Inches(6.7), Inches(2.0), Inches(2.8), Inches(2.8))
            notes_frame = notes_content.text_frame
            notes_frame.word_wrap = True
            
            for idx, note in enumerate(notes):
                if idx > 0:
                    p = notes_frame.add_paragraph()
                else:
                    p = notes_frame.paragraphs[0]
                p.text = f"• {note}"
                p.font.size = Pt(11)
                p.font.color.rgb = colors.DARK
                p.space_after = Pt(8)
        else:
            # Full width chart
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.1), width=Inches(9))
        
        print(f"   ✓ Slide created: {title}")
    else:
        print(f"   ✗ Chart not found: {chart_path}")
    
    return slide

def add_insights_slide(prs, insights_df):
    """Create insights slide with findings"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors.PRIMARY
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "EXECUTIVE INSIGHTS & RECOMMENDATIONS"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    
    # Insights cards
    y_pos = 1.2
    
    for idx, (_, insight) in enumerate(insights_df.iterrows()):
        if idx >= 4:  # Show max 4 insights
            break
        
        # Card background
        card = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.85))
        
        # Color based on type
        if insight['type'] == 'PRIMARY_DRIVER':
            card_color = colors.DANGER
        elif insight['type'] == 'OPPORTUNITY':
            card_color = colors.SUCCESS
        elif insight['type'] == 'TOP_FACTOR':
            card_color = colors.WARNING
        else:
            card_color = colors.INFO
        
        card.fill.solid()
        card.fill.fore_color.rgb = card_color
        card.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(8.6), Inches(0.3))
        title_box.text = insight['title']
        p = title_box.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors.WHITE
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.35), Inches(8.6), Inches(0.45))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        
        p = desc_frame.paragraphs[0]
        p.text = f"→ {insight['description']}"
        p.font.size = Pt(11)
        p.font.color.rgb = colors.WHITE
        
        # Action
        if pd.notna(insight['action']):
            p = desc_frame.add_paragraph()
            p.text = f"✓ {insight['action']}"
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        y_pos += 1.0
    
    print("   ✓ Insights slide created")
    return slide

def add_action_plan_slide(prs, category_impact):
    """Create action plan slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors.SUCCESS
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "RECOMMENDED ACTION PLAN"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors.WHITE
    
    # Priority actions based on top categories
    actions = [
        {
            'priority': '🔴 IMMEDIATE',
            'title': f'Address {category_impact.iloc[0]["Category"]}',
            'action': f'Launch task force to reduce {category_impact.iloc[0]["Category"].lower()} issues',
            'timeline': '30 days',
            'impact': 'High'
        },
        {
            'priority': '🟠 SHORT-TERM',
            'title': 'Improve Customer Experience',
            'action': 'Implement proactive support for at-risk customers',
            'timeline': '60 days',
            'impact': 'Medium-High'
        },
        {
            'priority': '🟡 MEDIUM-TERM',
            'title': 'Enhance Engagement Programs',
            'action': 'Develop loyalty rewards and referral incentives',
            'timeline': '90 days',
            'impact': 'Medium'
        },
        {
            'priority': '🟢 ONGOING',
            'title': 'Monitor & Optimize',
            'action': 'Weekly churn analytics dashboard and alerts',
            'timeline': 'Continuous',
            'impact': 'Preventive'
        }
    ]
    
    y_pos = 1.3
    
    for action in actions:
        # Action box
        box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = colors.LIGHT_BG
        box.line.color.rgb = colors.PRIMARY
        box.line.width = Pt(1)
        
        # Priority
        priority_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(2), Inches(0.25))
        priority_box.text = action['priority']
        p = priority_box.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = colors.DANGER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.35), Inches(5), Inches(0.3))
        title_box.text = action['title']
        p = title_box.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = colors.DARK
        
        # Details
        details = f"{action['action']} | Timeline: {action['timeline']} | Impact: {action['impact']}"
        details_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos + 0.1), Inches(6), Inches(0.25))
        details_box.text = details
        p = details_box.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = colors.SECONDARY
        
        y_pos += 0.85
    
    print("   ✓ Action plan slide created")
    return slide

def add_closing_slide(prs):
    """Create closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors.PRIMARY
    
    # Main message
    msg_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    msg_frame = msg_box.text_frame
    msg_frame.text = "LET'S REDUCE CHURN\nTOGETHER"
    
    for paragraph in msg_frame.paragraphs:
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.color.rgb = colors.WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Call to action
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.5))
    cta_box.text = "Questions? Let's discuss the action plan"
    p = cta_box.text_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = colors.INFO
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================================
# CREATE PRESENTATION
# ============================================================================

print("2. Creating executive presentation...")
print()

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

print("3. Adding slides...")
print()

# Slide 1: Title
add_title_slide(
    prs,
    "CUSTOMER CHURN\nCAUSALITY ANALYSIS",
    f"Understanding What Drives Customer Loss | {datetime.now().strftime('%B %Y')}"
)
print("   ✓ Slide 1: Title")

# Slide 2: Executive Summary
add_insights_slide(prs, insights)

# Slide 3: Hierarchical Causal Diagram
add_content_slide_with_chart(
    prs,
    "CAUSAL FACTOR CATEGORIES",
    f"{CHART_DIR}/causal_diagram_hierarchical.png",
    notes=[
        "Arrow thickness indicates causal strength",
        "Impact % shows contribution to churn",
        f"Top driver: {category_impact.iloc[0]['Category']}"
    ]
)

# Slide 4: Detailed Network
add_content_slide_with_chart(
    prs,
    "DETAILED CAUSAL NETWORK",
    f"{CHART_DIR}/causal_network_detailed.png"
)

# Slide 5: Waterfall
add_content_slide_with_chart(
    prs,
    "CHURN CONTRIBUTION ANALYSIS",
    f"{CHART_DIR}/churn_waterfall.png",
    notes=[
        "Shows cumulative impact",
        "Each bar adds to total churn",
        f"Top 3 factors = {feature_importance.head(3)['importance'].sum():.1%} impact"
    ]
)

# Slide 6: Sankey Flow
add_content_slide_with_chart(
    prs,
    "CHURN CAUSALITY FLOW",
    f"{CHART_DIR}/churn_sankey_flow.png",
    notes=[
        "Visualizes cause → effect flow",
        "Thicker flows = stronger causality",
        "Multiple factors compound"
    ]
)

# Slide 7: Action Plan
add_action_plan_slide(prs, category_impact)

# Slide 8: Closing
add_closing_slide(prs)
print("   ✓ Slide 8: Closing")

# ============================================================================
# SAVE PRESENTATION
# ============================================================================

print()
print("4. Saving presentation...")

output_path = f"{OUTPUT_DIR}/Churn_Causality_Executive_Report.pptx"
prs.save(output_path)

print(f"   ✓ Presentation saved: {output_path}")
print()

# ============================================================================
# SUMMARY
# ============================================================================

print("=" * 80)
print("EXECUTIVE PRESENTATION COMPLETE!")
print("=" * 80)
print()
print("✅ Presentation Details:")
print(f"   - Total Slides: 8")
print(f"   - Charts Included: 4 causal diagrams")
print(f"   - File Size: {os.path.getsize(output_path) / 1024:.1f} KB")
print()
print("📊 Content:")
print("   1. Title slide")
print("   2. Executive insights & findings")
print("   3. Hierarchical causal diagram")
print("   4. Detailed causal network")
print("   5. Churn waterfall analysis")
print("   6. Sankey flow diagram")
print("   7. Recommended action plan")
print("   8. Closing slide")
print()
print("💡 Management-Ready Features:")
print("   ✓ Clear causal relationships")
print("   ✓ Actionable insights")
print("   ✓ Prioritized recommendations")
print("   ✓ Visual explanations")
print("   ✓ Executive summary format")
print()
print("📥 Download Location:")
print(f"   {output_path}")
print()
print("=" * 80)
